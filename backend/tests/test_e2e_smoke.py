"""Smoke E2E：TestClient + FakeLLM 走完整 HTTP API 流程（規格第 9 節）。

與 test_api.py 的 happy path 差異：
- 用更完整的多頁大綱（cover / bullets / closing 三種 layout_hint）。
- 匯出後用 python-pptx 實際打開 pptx，斷言頁數、可尋文字內容都正確，
  而非只看 API 回應的 exported_count。
- 額外斷言 exports/ 目錄下的檔案確實落地在磁碟上（不是只存在於回應）。
"""

import io

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from app.api.deps import get_llm, get_projects_root
from app.main import app
from tests.conftest import FakeLLM

# 三頁大綱：封面 / 條列重點 / 結語，各自代表一種 layout_hint。
_OUTLINE_JSON = """```json
{"slides": [
  {"index": 0, "title": "封面：2026 產品發表會", "bullets": ["驚豔登場"], "layout_hint": "cover", "assets": []},
  {"index": 1, "title": "三大重點", "bullets": ["效能提升三倍", "全新使用者介面", "價格更親民"], "layout_hint": "bullets", "assets": []},
  {"index": 2, "title": "結語：感謝參與", "bullets": ["敬請期待下次更新"], "layout_hint": "closing", "assets": []}
]}
```"""

# 每頁預錄一段可辨識的繁中文字，供匯出後在 pptx 文字內容中尋找比對。
_SLIDE_TEXTS = ["封面：2026 產品發表會", "三大重點：效能提升三倍", "結語：感謝參與，敬請期待下次更新"]


def _svg(text: str) -> str:
    return (
        "```svg\n"
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">'
        '<rect x="0" y="0" width="1280" height="720" fill="#1a2b3c"/>'
        f'<text x="40" y="100" font-size="24">{text}</text>'
        "</svg>\n```"
    )


@pytest.fixture
def projects_root(tmp_path):
    root = tmp_path / "projects"
    root.mkdir()
    return root


@pytest.fixture
def client(projects_root):
    app.dependency_overrides[get_projects_root] = lambda: projects_root
    with TestClient(app) as c:
        yield c
    app.dependency_overrides.pop(get_projects_root, None)
    app.dependency_overrides.pop(get_llm, None)


def _override_llm(responses: list[str]) -> FakeLLM:
    fake = FakeLLM(responses)
    app.dependency_overrides[get_llm] = lambda: fake
    return fake


def _poll_progress(client, project_id, timeout=5.0):
    import time

    deadline = time.time() + timeout
    while time.time() < deadline:
        resp = client.get(f"/api/projects/{project_id}/progress")
        assert resp.status_code == 200
        data = resp.json()
        if data["stage"] == "generated":
            return data
        time.sleep(0.02)
    raise AssertionError("generation did not complete in time")


_SAMPLE_MD = """# 2026 產品發表會

## 三大重點
- 效能提升三倍
- 全新使用者介面
- 價格更親民

## 結語
感謝參與，敬請期待下次更新。
"""


def test_e2e_smoke_upload_to_export(client, projects_root):
    """上傳 sample.md -> 選風格 -> outline -> generate -> 輪詢 -> 匯出。

    匯出結果用 python-pptx 打開驗證：頁數正確、文字可尋、檔案落地。
    """
    # 1. 建立專案
    resp = client.post("/api/projects", json={"name": "E2E 煙霧測試"})
    assert resp.status_code == 200, resp.text
    project_id = resp.json()["id"]

    # 2. 上傳 sample.md
    files = {"files": ("sample.md", io.BytesIO(_SAMPLE_MD.encode("utf-8")), "text/markdown")}
    resp = client.post(f"/api/projects/{project_id}/upload", files=files)
    assert resp.status_code == 200, resp.text
    assert resp.json()["results"][0]["success"] is True

    # 3. 選風格
    styles = client.get("/api/styles").json()
    assert len(styles["styles"]) >= 1
    assert len(styles["palettes"]) >= 1
    resp = client.post(
        f"/api/projects/{project_id}/style",
        json={
            "style_id": styles["styles"][0]["id"],
            "palette_id": styles["palettes"][0]["id"],
        },
    )
    assert resp.status_code == 200, resp.text

    # 4. 產生 outline（FakeLLM 預錄 3 頁）
    _override_llm([_OUTLINE_JSON])
    resp = client.post(f"/api/projects/{project_id}/outline")
    assert resp.status_code == 200, resp.text
    outline = resp.json()
    assert len(outline["slides"]) == 3
    assert [s["layout_hint"] for s in outline["slides"]] == ["cover", "bullets", "closing"]

    # 5. 觸發生成（背景任務；FakeLLM 依序回傳 3 頁 SVG）
    _override_llm([_svg(t) for t in _SLIDE_TEXTS])
    resp = client.post(f"/api/projects/{project_id}/generate")
    assert resp.status_code == 202, resp.text

    # 6. 輪詢 progress 至完成
    progress = _poll_progress(client, project_id)
    assert progress["stage"] == "generated"
    assert len(progress["slides"]) == 3
    assert all(s["status"] == "generated" for s in progress["slides"])

    # 7. 逐頁 SVG 可取得且內容含預錄文字
    for i, expected_text in enumerate(_SLIDE_TEXTS):
        resp = client.get(f"/api/projects/{project_id}/slides/{i}.svg")
        assert resp.status_code == 200
        assert resp.headers["content-type"].startswith("image/svg+xml")
        assert expected_text in resp.text

    # 8. 匯出
    resp = client.post(f"/api/projects/{project_id}/export")
    assert resp.status_code == 200, resp.text
    export_result = resp.json()
    assert export_result["exported_count"] == 3
    assert export_result["skipped_count"] == 0
    download_url = export_result["download_url"]
    assert download_url

    # 9. 下載 pptx
    resp = client.get(download_url)
    assert resp.status_code == 200
    assert resp.headers["content-type"] in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/octet-stream",
    )
    pptx_bytes = resp.content
    assert len(pptx_bytes) > 0

    # 10. exports/ 目錄下的檔案確實落地在磁碟上（不是只存在於 HTTP 回應）
    exports_dir = projects_root / project_id / "exports"
    exported_files = list(exports_dir.glob("*.pptx"))
    assert len(exported_files) == 1, f"預期 exports/ 下恰有 1 個 pptx，實際：{exported_files}"
    filename = download_url.rsplit("/", 1)[-1]
    assert exported_files[0].name == filename
    assert exported_files[0].read_bytes() == pptx_bytes

    # 11. 用 python-pptx 打開匯出的檔案，斷言頁數與文字內容
    prs = Presentation(exported_files[0])
    assert len(prs.slides) == 3

    all_slide_texts = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        all_slide_texts.append("\n".join(texts))

    # 每一頁的預錄文字都要能在對應頁的 shape 文字中尋得。
    for i, expected_text in enumerate(_SLIDE_TEXTS):
        assert expected_text in all_slide_texts[i], (
            f"第 {i} 頁文字「{expected_text}」未在匯出的 pptx 中找到，"
            f"實際內容：{all_slide_texts[i]!r}"
        )
