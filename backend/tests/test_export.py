"""export_pptx() 測試：vendored svg_to_pptx 轉換與元素級降級。"""

import re

import pytest
from pptx import Presentation

from app.export.pptx_export import ExportError, export_pptx
from app.store.project import create_project

_RECT_SVG = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
<rect x="10" y="10" width="200" height="100" fill="#336699"/>
<text x="20" y="60" font-size="24">{text}</text>
</svg>"""

_INVALID_SVG = "<svg xmlns=\"http://www.w3.org/2000/svg\" viewBox=\"0 0 1280 720\"><rect><<<broken"


def _write_slide(project, index, text):
    svg_dir = project.path / "svg_output"
    svg_dir.mkdir(parents=True, exist_ok=True)
    (svg_dir / f"slide_{index:03d}.svg").write_text(
        _RECT_SVG.format(text=text), encoding="utf-8"
    )


def _make_project_with_slides(tmp_path, texts: list[str]):
    project = create_project(tmp_path, "匯出測試")
    for i, text in enumerate(texts):
        _write_slide(project, i, text)
        project.set_slide_status(i, "generated")
    project.save()
    return project


def test_export_produces_editable_pptx(tmp_path):
    project = _make_project_with_slides(tmp_path, ["第一頁標題", "第二頁標題"])

    result = export_pptx(project)

    assert result.output_path.is_file()
    assert result.exported_count == 2
    assert result.skipped_count == 0
    assert result.warnings == []

    prs = Presentation(str(result.output_path))
    assert len(prs.slides) == 2

    first_slide_texts = [
        shape.text_frame.text
        for shape in prs.slides[0].shapes
        if shape.has_text_frame
    ]
    assert any("第一頁標題" in t for t in first_slide_texts)


def test_failed_slide_skipped_with_warning(tmp_path):
    project = create_project(tmp_path, "部分失敗")
    _write_slide(project, 0, "合法頁")
    project.set_slide_status(0, "generated")
    # index 1 從未生成，狀態停在 pending -> 不應被匯出
    project.save()

    result = export_pptx(project)

    assert result.exported_count == 1
    prs = Presentation(str(result.output_path))
    assert len(prs.slides) == 1


def test_all_pages_missing_raises(tmp_path):
    project = create_project(tmp_path, "全無生成頁")
    project.save()

    with pytest.raises(ExportError):
        export_pptx(project)


def test_conversion_failure_falls_back_to_rasterized_image(tmp_path, monkeypatch):
    project = _make_project_with_slides(tmp_path, ["正常頁"])
    # 追加一頁內容合法但會讓 vendor 轉換器丟例外的 SVG（用 monkeypatch 模擬單頁轉換失敗，
    # 而不是真的塞壞掉的 XML，因為壞 XML 在生成階段就會被 check_svg 擋下；
    # 這裡驗證的是匯出層對「轉換器本身丟例外」的降級路徑）。
    import app.export.pptx_export as pptx_export_module

    original = pptx_export_module.convert_svg_to_slide_shapes

    call_count = {"n": 0}

    def _boom(svg_path, *args, **kwargs):
        call_count["n"] += 1
        if call_count["n"] == 1:
            raise RuntimeError("模擬轉換失敗")
        return original(svg_path, *args, **kwargs)

    monkeypatch.setattr(pptx_export_module, "convert_svg_to_slide_shapes", _boom)

    result = export_pptx(project)

    assert result.exported_count == 1
    assert result.skipped_count == 0
    assert len(result.warnings) == 1
    assert "slide_000" in result.warnings[0] or "第 1 頁" in result.warnings[0]

    prs = Presentation(str(result.output_path))
    assert len(prs.slides) == 1
    # 降級頁應含至少一個圖片 shape（rasterized PNG 塞成整頁圖片）
    assert any(shape.shape_type == 13 for shape in prs.slides[0].shapes)  # 13 = PICTURE


def test_output_filename_timestamp_and_sanitized_name(tmp_path):
    project = _make_project_with_slides(tmp_path, ["頁"])
    project.data["name"] = "月報 / 測試 draft"
    project.save()

    result = export_pptx(project)

    filename = result.output_path.name
    assert filename.endswith(".pptx")
    assert re.search(r"_\d{8}_\d{6}\.pptx$", filename)
    assert "/" not in filename.removesuffix(".pptx")
    assert " " not in filename
    assert result.output_path.parent == project.path / "exports"


def test_image_href_escaping_project_root_is_skipped_with_warning(tmp_path):
    # 在 project.path 之外放一個「真實存在」的檔案：候選路徑 resolve 後
    # 確實存在，防護必須靠 relative_to(project_root) 的 ValueError 分支
    # 才擋得下來。若有人刪掉該段防護，這個測試會確定性地轉紅（不像指向
    # 不存在路徑的 href，那只會命中「候選不存在」分支，防護刪了照樣綠）。
    outside_file = tmp_path / "outside.txt"
    outside_file.write_bytes(b"\x89PNG fake-but-real-file")

    project = create_project(tmp_path, "href逃逸")
    svg_dir = project.path / "svg_output"
    svg_dir.mkdir(parents=True, exist_ok=True)
    # 佈局為 tmp_path/<project_id>/svg_output/，因此從 svg_output/ 走
    # ../../outside.txt 剛好落在 tmp_path/outside.txt（project.path 之外）。
    escaping_svg = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
<image href="../../outside.txt" x="0" y="0" width="100" height="100"/>
<text x="20" y="60" font-size="24">安全測試頁</text>
</svg>"""
    (svg_dir / "slide_000.svg").write_text(escaping_svg, encoding="utf-8")
    project.set_slide_status(0, "generated")
    project.save()

    # 前置驗證：這個 href 的候選路徑確實存在且在專案外——確保測試打到
    # relative_to 分支，而不是「檔案不存在」分支。
    resolved = (svg_dir / "../../outside.txt").resolve()
    assert resolved.is_file()
    assert not resolved.is_relative_to(project.path.resolve())

    result = export_pptx(project)

    # 逃逸的 image 應被略過並記警告；頁面仍應成功匯出（純降級該圖片，而非整頁）
    # 而不是讓匯出動作讀取到專案目錄之外的檔案內容。
    assert result.exported_count == 1
    assert len(result.warnings) == 1
    assert "outside.txt" in result.warnings[0]
    prs = Presentation(str(result.output_path))
    assert len(prs.slides) == 1


def test_vendor_private_symbols_still_importable():
    """vendor 升級相容性 smoke test：adapter 依賴的 vendor 符號若被改名，
    這裡會直接 ImportError／AssertionError 轉紅，訊息明確指向 vendor 介面變動。"""
    from svg_to_pptx.drawingml.converter import convert_svg_to_slide_shapes
    from svg_to_pptx.pptx_package.builder import (
        _add_default_content_type,
        _content_type_for_extension,
        _create_writable_work_dir,
    )

    assert callable(convert_svg_to_slide_shapes)
    assert callable(_add_default_content_type)
    assert callable(_content_type_for_extension)
    assert callable(_create_writable_work_dir)
