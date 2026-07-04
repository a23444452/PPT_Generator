"""匯出 PPTX：vendor svg_to_pptx 逐頁轉換 + 元素級降級。

Vendor 進入點探查（Task 9 Step 2 記錄）：
    vendor/svg_to_pptx/__init__.py 的公開 API 是
        create_pptx_with_native_svg(svg_files: list[Path], output_path: Path, ...) -> bool
    這個函式接受單純的 SVG 檔案清單與輸出路徑，不依賴 ppt-master 的整個
    project 目錄結構（那是 CLI 層 `pptx_package/cli.py` 才做的事：掃描
    svg_output/、讀 notes/、animations.json 等）。

    但 create_pptx_with_native_svg 本身「單頁失敗即整檔失敗」：native
    shapes 模式下，逐頁迴圈裡任何例外都會直接 raise（builder.py 第
    1144-1148 行 `if use_native_shapes: raise`），沒有「這頁降級、其餘頁
    照常」的機制。這與本任務要求的「單頁轉換失敗→該頁降級為 rasterized
    圖片、繼續匯出其餘頁」衝突，且不能修改 vendor 內部檔案，因此本模組
    複用 create_pptx_with_native_svg 內部同一套 pptx 組裝手法（python-pptx
    建立 N 頁空白版面 → 存檔 → unzip → 逐頁寫入 slideN.xml/rels →
    補 [Content_Types].xml → rezip），但把「呼叫 vendor 的
    drawingml.converter.convert_svg_to_slide_shapes()」這一步包在
    try/except 內：成功則寫入 vendor 產生的原生 DrawingML slide XML，
    失敗則改用 PyMuPDF 把該頁 SVG rasterize 成 PNG、寫一份「整頁圖片」
    slide XML 頂替，並記警告——組裝骨架是本模組自己的程式碼，但真正的
    SVG→DrawingML 轉換邏輯、以及 zip 組裝用的小型 helper（
    _create_writable_work_dir、_add_default_content_type、
    _content_type_for_extension）都直接 import 自 vendor 的
    pptx_package.builder，未複製或修改其原始碼。

    頁面尺寸：MVP 固定 16:9、1280×720px viewBox（見
    app/generation/quality.py 的 EXPECTED_VIEWBOX），對應
    slide_width/height = Inches(13.333) x Inches(7.5)。

Vendor 耦合與 adapter 範圍（不修改 vendor 內部檔案）：
    svg_to_pptx 的 drawingml/converter.py、elements.py 以「裸模組名」
    （非相對匯入）匯入 resource_paths、console_encoding；
    tspan_flattener.py 進一步 lazy-import svg_finalize.flatten_tspan
    （svg_finalize 也只依賴 console_encoding，無其他第三方套件）。
    這些都是 vendor 原始碼中既有的硬相依，因此把這三個 sibling 模組
    （console_encoding.py、resource_paths.py、svg_finalize/）原封不動
    複製到 vendor/ 目錄下、與 svg_to_pptx/ 同層，並在載入前把 vendor/
    根目錄插入 sys.path，讓這些裸匯入能解析到——不改 vendor 內任何檔案，
    日後升級只需整包替換 vendor/svg_to_pptx/ 與這三個 sibling 模組。

    dimensions.py 匯入的 project_utils / config 有 try/except 保底
    （匯入失敗則用內建預設 CANVAS_FORMATS），不必額外 vendor。

安全注意（href 逃逸防護）：
    resource_paths.resolve_external_image_reference() 只檢查候選路徑
    是否存在，不檢查是否逃出 project 根目錄——如果專案外剛好有同名檔案，
    <image href="../../..."> 可能讀到專案目錄外的檔案。quality.py 的
    href 檢查只在生成當下做格式白名單（assets/ 前綴或 data URI），並非
    路徑安全邊界，匯出這個實際讀檔的時間點才需要真正防護。因此本模組在
    呼叫 vendor 轉換器之前，會先掃描每頁 SVG 的 <image> href，只要是
    非 data: 的 href，一律 resolve 後確認結果路徑在 project.path 之下；
    逃逸或不存在則從 SVG 樹中移除該 image 節點並記警告，避免 vendor
    讀取到專案目錄以外的檔案。
"""

from __future__ import annotations

import re
import shutil
import sys
import zipfile
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from urllib.parse import unquote, urlsplit
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.util import Inches

from app.store.project import Project

_VENDOR_ROOT = Path(__file__).resolve().parents[3] / "vendor"
if str(_VENDOR_ROOT) not in sys.path:
    sys.path.insert(0, str(_VENDOR_ROOT))

from svg_to_pptx.drawingml.converter import convert_svg_to_slide_shapes  # noqa: E402
from svg_to_pptx.pptx_package.builder import (  # noqa: E402
    _add_default_content_type,
    _content_type_for_extension,
    _create_writable_work_dir,
)

_SVG_NS = "http://www.w3.org/2000/svg"
_XLINK_NS = "http://www.w3.org/1999/xlink"

_SLIDE_WIDTH_EMU = Inches(13.333)
_SLIDE_HEIGHT_EMU = Inches(7.5)
_TIMESTAMP_FORMAT = "%Y%m%d_%H%M%S"
_SANITIZE_RE = re.compile(r"[\\/\s]+")


class ExportError(Exception):
    """匯出失敗（訊息對使用者友善，不洩漏內部細節）。"""


@dataclass
class ExportResult:
    output_path: Path
    warnings: list[str] = field(default_factory=list)
    exported_count: int = 0
    skipped_count: int = 0


def export_pptx(project: Project) -> ExportResult:
    """將 project 內所有 status=="generated" 的 SVG 依序轉為一份 pptx。

    輸出檔名：exports/<name>_<YYYYMMDD_HHMMSS>.pptx（name 已清理空白／斜線）。
    單頁轉換失敗會降級為整頁 rasterized PNG 圖片並記警告，不中止其餘頁面；
    全部頁面都無法匯出（無 generated 頁，或全部轉換失敗）才拋出 ExportError。
    """
    svg_paths = _collect_generated_svg_paths(project)
    if not svg_paths:
        raise ExportError("沒有已生成的頁面可供匯出，請先完成投影片生成。")

    output_path = _build_output_path(project)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    temp_dir = _create_writable_work_dir(output_path)
    try:
        prs = Presentation()
        prs.slide_width = _SLIDE_WIDTH_EMU
        prs.slide_height = _SLIDE_HEIGHT_EMU
        blank_layout = prs.slide_layouts[6]
        for _ in svg_paths:
            prs.slides.add_slide(blank_layout)

        base_pptx = temp_dir / "base.pptx"
        prs.save(str(base_pptx))

        extract_dir = temp_dir / "pptx_content"
        with zipfile.ZipFile(base_pptx, "r") as zf:
            zf.extractall(extract_dir)

        media_dir = extract_dir / "ppt" / "media"
        media_dir.mkdir(exist_ok=True)

        export_warnings: list[str] = []
        image_exts_used: set[str] = set()
        slides_written = 0
        slides_to_remove: list[int] = []

        for i, svg_path in enumerate(svg_paths, start=1):
            try:
                href_warnings = _write_native_slide(
                    project, extract_dir, media_dir, svg_path, i, image_exts_used
                )
                for href in href_warnings:
                    export_warnings.append(
                        f"第 {i} 頁（{svg_path.name}）圖片參照「{href}」"
                        "超出專案目錄範圍，已略過該圖片。"
                    )
                slides_written += 1
            except Exception as exc:  # noqa: BLE001 — 單頁降級，需捕捉任何轉換例外
                try:
                    _write_rasterized_fallback_slide(
                        extract_dir, media_dir, svg_path, i, image_exts_used
                    )
                    export_warnings.append(
                        f"第 {i} 頁（{svg_path.name}）原生轉換失敗，已降級為圖片：{exc}"
                    )
                    slides_written += 1
                except Exception as fallback_exc:  # noqa: BLE001
                    export_warnings.append(
                        f"第 {i} 頁（{svg_path.name}）轉換與降級都失敗，已略過："
                        f"{fallback_exc}"
                    )
                    slides_to_remove.append(i)

        if slides_written == 0:
            raise ExportError("所有頁面轉換皆失敗，無法產生匯出檔案。")

        if slides_to_remove:
            _drop_failed_slides(extract_dir, slides_to_remove, len(svg_paths))

        _update_content_types(extract_dir, image_exts_used)

        with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for file_path in extract_dir.rglob("*"):
                if file_path.is_file():
                    zf.write(file_path, file_path.relative_to(extract_dir))

        return ExportResult(
            output_path=output_path,
            warnings=export_warnings,
            exported_count=slides_written,
            skipped_count=len(svg_paths) - slides_written,
        )
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


def _collect_generated_svg_paths(project: Project) -> list[Path]:
    svg_dir = project.path / "svg_output"
    paths: list[Path] = []
    for slide in project.data.get("slides", []):
        if slide.get("status") != "generated":
            continue
        svg_path = svg_dir / f"slide_{slide['index']:03d}.svg"
        if svg_path.is_file():
            paths.append(svg_path)
    return paths


def _build_output_path(project: Project) -> Path:
    timestamp = datetime.now().strftime(_TIMESTAMP_FORMAT)
    safe_name = _SANITIZE_RE.sub("_", project.data.get("name", project.id)).strip("_")
    if not safe_name:
        safe_name = project.id
    return project.path / "exports" / f"{safe_name}_{timestamp}.pptx"


def _write_native_slide(
    project: Project,
    extract_dir: Path,
    media_dir: Path,
    svg_path: Path,
    slide_num: int,
    image_exts_used: set[str],
) -> list[str]:
    """呼叫 vendor 轉換器把單一 SVG 轉為原生 DrawingML shapes，寫入 slide part。

    回傳本頁被移除的逃逸 href 清單（供呼叫端記警告；空清單表示沒有）。
    """
    guarded_path, href_warnings = _guard_image_hrefs(project, svg_path)
    try:
        (
            slide_xml,
            media_files,
            rel_entries,
            _anim_targets,
            _package_files,
            _content_type_overrides,
        ) = convert_svg_to_slide_shapes(guarded_path, slide_num=slide_num, verbose=False)
    finally:
        if guarded_path != svg_path:
            guarded_path.unlink(missing_ok=True)

    slide_xml_path = extract_dir / "ppt" / "slides" / f"slide{slide_num}.xml"
    slide_xml_path.write_text(slide_xml, encoding="utf-8")

    for media_name, media_data in media_files.items():
        with open(media_dir / media_name, "wb") as f:
            f.write(media_data)
        ext = media_name.rsplit(".", 1)[-1].lower()
        image_exts_used.add(ext)

    rels_dir = extract_dir / "ppt" / "slides" / "_rels"
    rels_dir.mkdir(exist_ok=True)
    rels_path = rels_dir / f"slide{slide_num}.xml.rels"

    extra_rels = ""
    for rel in rel_entries:
        extra_rels += (
            f'\n  <Relationship Id="{rel["id"]}" '
            f'Type="{rel["type"]}" Target="{rel["target"]}"/>'
        )

    rels_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1"
                Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
                Target="../slideLayouts/slideLayout1.xml"/>{extra_rels}
</Relationships>'''
    rels_path.write_text(rels_xml, encoding="utf-8")

    return href_warnings


def _write_rasterized_fallback_slide(
    extract_dir: Path,
    media_dir: Path,
    svg_path: Path,
    slide_num: int,
    image_exts_used: set[str],
) -> None:
    """降級：整頁 rasterize 成 PNG，寫一份純圖片 slide XML 頂替原生轉換。"""
    import fitz

    svg_bytes = svg_path.read_bytes()
    doc = fitz.open(stream=svg_bytes, filetype="svg")
    try:
        page = doc[0]
        pix = page.get_pixmap(dpi=150)
        png_bytes = pix.tobytes("png")
    finally:
        doc.close()

    media_name = f"fallback_{slide_num}.png"
    (media_dir / media_name).write_bytes(png_bytes)
    image_exts_used.add("png")

    width_emu = int(_SLIDE_WIDTH_EMU)
    height_emu = int(_SLIDE_HEIGHT_EMU)
    slide_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
       xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
<p:cSld>
<p:spTree>
<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/>
<a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
<p:pic>
<p:nvPicPr>
<p:cNvPr id="2" name="fallback-image"/>
<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
<p:nvPr/>
</p:nvPicPr>
<p:blipFill><a:blip r:embed="rId2"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>
<p:spPr>
<a:xfrm><a:off x="0" y="0"/><a:ext cx="{width_emu}" cy="{height_emu}"/></a:xfrm>
<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
</p:spPr>
</p:pic>
</p:spTree>
</p:cSld>
<p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sld>'''

    slide_xml_path = extract_dir / "ppt" / "slides" / f"slide{slide_num}.xml"
    slide_xml_path.write_text(slide_xml, encoding="utf-8")

    rels_dir = extract_dir / "ppt" / "slides" / "_rels"
    rels_dir.mkdir(exist_ok=True)
    rels_path = rels_dir / f"slide{slide_num}.xml.rels"
    rels_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1"
                Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
                Target="../slideLayouts/slideLayout1.xml"/>
  <Relationship Id="rId2"
                Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
                Target="../media/{media_name}"/>
</Relationships>'''
    rels_path.write_text(rels_xml, encoding="utf-8")


def _drop_failed_slides(extract_dir: Path, failed_indices: list[int], total: int) -> None:
    """從 presentation.xml / rels / [Content_Types].xml 移除完全失敗的頁。

    python-pptx 建立 base pptx 時已經替每一頁配好 sldIdLst／關聯與
    slideN.xml part 引用；沒有寫入對應 slideN.xml 的頁若留著會讓
    PowerPoint 視為套件損毀，因此連同 presentation part 的引用一併清除。
    """
    ppt_dir = extract_dir / "ppt"
    presentation_path = ppt_dir / "presentation.xml"
    presentation_rels_path = ppt_dir / "_rels" / "presentation.xml.rels"

    presentation_xml = presentation_path.read_text(encoding="utf-8")
    rels_xml = presentation_rels_path.read_text(encoding="utf-8")

    # 找出失敗頁對應的 slideN.xml part 在 rels 裡的 rId，藉此在 sldIdLst
    # 與 rels 內同步移除；base pptx 是循序建立的，slide part 命名固定為
    # slide{1-based index}.xml。
    for index in failed_indices:
        target = f"slides/slide{index}.xml"
        match = re.search(
            rf'<Relationship Id="(rId\d+)"[^>]*Target="{re.escape(target)}"[^>]*/>',
            rels_xml,
        )
        if not match:
            continue
        rid = match.group(1)
        rels_xml = re.sub(
            rf'\s*<Relationship Id="{rid}"[^>]*/>', "", rels_xml
        )
        presentation_xml = re.sub(
            rf'\s*<p:sldId[^>]*r:id="{rid}"[^>]*/>', "", presentation_xml
        )
        slide_path = ppt_dir / "slides" / f"slide{index}.xml"
        slide_path.unlink(missing_ok=True)
        slide_rels_path = ppt_dir / "slides" / "_rels" / f"slide{index}.xml.rels"
        slide_rels_path.unlink(missing_ok=True)

    presentation_path.write_text(presentation_xml, encoding="utf-8")
    presentation_rels_path.write_text(rels_xml, encoding="utf-8")


def _update_content_types(extract_dir: Path, image_exts_used: set[str]) -> None:
    content_types_path = extract_dir / "[Content_Types].xml"
    content_types = content_types_path.read_text(encoding="utf-8")
    for ext in sorted(image_exts_used):
        content_types = _add_default_content_type(
            content_types, ext, _content_type_for_extension(ext)
        )
    content_types_path.write_text(content_types, encoding="utf-8")


def _guard_image_hrefs(project: Project, svg_path: Path) -> tuple[Path, list[str]]:
    """掃描 SVG 的 <image> href，移除任何會逃出 project.path 的參照。

    只處理非 data: 的 href：resolve 候選路徑並確認其在 project.path 之下；
    逃逸或找不到檔案則整個移除該 <image> 節點（避免 vendor 轉換器嘗試讀取
    專案目錄外的檔案）。若沒有任何節點需要移除，回傳原始路徑；否則寫一份
    清理過的暫存 SVG 供轉換器使用（呼叫端用完需自行清理暫存檔）。
    """
    original_bytes = svg_path.read_bytes()
    try:
        ET.register_namespace("", _SVG_NS)
        tree = ET.ElementTree(ET.fromstring(original_bytes))
    except ET.ParseError:
        return svg_path, []

    root = tree.getroot()
    project_root = project.path.resolve()
    svg_dir = svg_path.parent.resolve()
    warnings_found: list[str] = []

    parent_map = {child: parent for parent in root.iter() for child in parent}
    for image_elem in list(root.iter(f"{{{_SVG_NS}}}image")):
        href = image_elem.get("href") or image_elem.get(f"{{{_XLINK_NS}}}href")
        if href is None or href.strip().startswith("data:"):
            continue
        if not _href_resolves_within_root(svg_dir, project_root, href):
            parent = parent_map.get(image_elem)
            if parent is not None:
                parent.remove(image_elem)
                warnings_found.append(href)

    if not warnings_found:
        return svg_path, []

    tmp_path = svg_path.with_name(svg_path.stem + ".guarded.svg")
    tree.write(tmp_path, xml_declaration=True)
    return tmp_path, warnings_found


def _href_resolves_within_root(svg_dir: Path, project_root: Path, href: str) -> bool:
    parsed = urlsplit(href)
    if parsed.scheme and parsed.scheme not in {"file", ""}:
        return False
    decoded = unquote(
        parsed.path if parsed.scheme else href.split("?", 1)[0].split("#", 1)[0]
    )
    candidates = [
        svg_dir / decoded,
        project_root / decoded,
        project_root / "assets" / decoded,
    ]
    for candidate in candidates:
        try:
            resolved = candidate.resolve()
        except OSError:
            continue
        if not resolved.is_file():
            continue
        try:
            resolved.relative_to(project_root)
        except ValueError:
            continue
        return True
    return False
